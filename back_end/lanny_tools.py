import csv
import json
import logging
import os
import re
import uuid
from datetime import datetime
from typing import List

import chromadb
from PyPDF2 import PdfReader
from bs4 import BeautifulSoup
from chromadb.utils.embedding_functions.openai_embedding_function import OpenAIEmbeddingFunction
from openai import OpenAI
from dotenv import load_dotenv

from mysql_tools import DataBaseManage
from constant import LLM_CLOUD_URL, LLM_MODEL, MESSAGES


class Chatbot:
    def __init__(self):
        # 加载 .env 文件
        load_dotenv()
        api_key = os.getenv('QIANWEN_API_KEY')
        # 设置本地服务器地址
        # self.client = OpenAI(base_url=f"{LLM_BASE_URL}/v1", api_key="lm_studio")
        self.client = OpenAI(base_url=f"{LLM_CLOUD_URL}", api_key=api_key)

    """对话函数，没有实现上下文"""

    def chatbot(self, messages, query):
        # context = get_context()
        # messages.extend(context)
        # 进行rag
        embedding_func = EmbeddingFunc()
        res = embedding_func.search_in_vector_db(query)

        # logging.info(f"--------{messages}")
        # 构建提示词
        prompt = f"""
        # 指导原则
        在回答以下问题时，请遵循以下逻辑：
        1. **优先使用检索到的内容**：如果检索到的信息与问题直接相关，请根据这些信息生成回答。
        2. **补充 LLM 的知识**：如果检索到的内容与问题无关，或者不足以回答问题，请结合你自己的知识进行回答。

        # 检索到的内容
        以下是检索到的相关信息：
        - {res}

        # 用户问题
        {query}

        # 回答
        根据上述逻辑，请生成回答：

        示例：
            -哟，你问啥子AI哦！不就是那种能自己动脑筋的机器嘛！就像电影里的机器人，能自己做决定，还能帮你干好多事。厉害得很哦，但有时候也像个智障，哈哈！
            -哎哟喂，你们这些外地人，问起AI能干啥子，老子今天就给你们摆摆清楚！你们晓得不，我们四川现在连挑二荆条辣椒都要用AI来帮忙选，搞得比我们这些老四川人还专业！
            -哟，摆啥子龙门阵！你问啥子AI Agent，老子给你摆一摆！这玩意儿就是个高科技二哈，专门在电脑里面自己做主的智能体！它能自己看、自己听、自己想，还能自己动手干事情，比你家那只只会汪汪叫的狗子强多了！
        
        """
        messages.append({"role": "user", "content": prompt})
        # messages.append({"role": "user", "content": query})
        response = self.client.chat.completions.create(
            model=LLM_MODEL,
            # model="qwen2.5-7b-instruct-mlx",
            messages=messages,
            temperature=0.7,
            stream=True,
            max_tokens=512,
        )

        def generate_response():
            assistant_response = ""
            response_chunks = []
            chunk_id = 1

            for chunk in response:
                if chunk.choices[0].delta.content:
                    content = chunk.choices[0].delta.content
                    assistant_response += content
                    # print(content, end='')
                    response_chunks.append({"chunk_id": chunk_id, "content": content})
                    chunk_id += 1
                    # logging.info(content)
                    yield content  # 流式返回每个 chunk 的内容
            messages.append({"role": "user", "content": query})
            messages.append({"role": "assistant", "content": assistant_response})

            # messages.append({"role": "user", "content": assistant_response})
            # 构造结构化输出
            structured_output = {
                "user_input": query,
                "model_response": assistant_response,
                "timestamp": datetime.utcnow().isoformat(),
                "messages": messages,
            }
            logging.info(structured_output)
            # 在流式返回的最后，返回结构化输出
            # yield f"{structured_output}"
            # insert_context()

        return generate_response()

    """对话函数，实现上下文存储，历史记录读取"""

    def chatbots(self, user_id, chatbot_id, query):
        messages = []
        # 获取上下文
        dataBM = DataBaseManage()
        context = dataBM.get_context(user_id, chatbot_id)
        logging.info(f"上下文--------{context}")
        # 将获取到的近10条上下文加入到messages中
        messages.extend(context)
        messages.extend(MESSAGES)
        messages.append({"role": "user", "content": query})
        # 进行rag
        embedding_func = EmbeddingFunc()

        res = embedding_func.search_in_vector_db(query, user_id, chatbot_id)
        # logging.info(f"查询结果--------{res}")

        # 根据rag结果   构建提示词
        prompt = f"""
        # 指导原则
        在回答以下问题时，请遵循以下逻辑：
        1. **优先使用检索到的内容**：如果检索到的信息与问题直接相关，请根据这些信息生成回答。
        2. **补充 LLM 的知识**：如果检索到的内容与问题无关，或者不足以回答问题，请不要使用检索到的内容进行回答！！！请仅你自己的知识进行回答。

        # 检索到的内容
        以下是检索到的相关信息：
        - {res}

        # 用户问题
        {query}

        # 回答
        根据上述逻辑，请生成回答：

        示例：
            -哟，你问啥子AI哦！不就是那种能自己动脑筋的机器嘛！就像电影里的机器人，能自己做决定，还能帮你干好多事。厉害得很哦，但有时候也像个智障，哈哈！
            -哎哟喂，你们这些外地人，问起AI能干啥子，老子今天就给你们摆摆清楚！你们晓得不，我们四川现在连挑二荆条辣椒都要用AI来帮忙选，搞得比我们这些老四川人还专业！
            -哟，摆啥子龙门阵！你问啥子AI Agent，老子给你摆一摆！这玩意儿就是个高科技二哈，专门在电脑里面自己做主的智能体！它能自己看、自己听、自己想，还能自己动手干事情，比你家那只只会汪汪叫的狗子强多了！

        """
        messages.append({"role": "user", "content": prompt.strip()})  # 添加strip()去除多余空白
        # logging.info(f'{messages},,{type(messages)}')
        try:
            logging.info(f",{type(messages)}")
            response = self.client.chat.completions.create(
                model=LLM_MODEL,
                messages=messages,
                temperature=0.7,
                stream=True,
                max_tokens=512,
            )
        except Exception as e:
            logging.error(f"发生错误：{e}")
            return {'status': False, 'message': str(e)}

        def generate_response():
            assistant_response = ""
            response_chunks = []
            chunk_id = 1

            for chunk in response:
                if chunk.choices[0].delta.content:
                    content = chunk.choices[0].delta.content
                    assistant_response += content
                    # print(content, end='')
                    response_chunks.append({"chunk_id": chunk_id, "content": content})
                    chunk_id += 1
                    # logging.info(content)
                    yield content  # 流式返回每个 chunk 的内容
            messages.append({"role": "user", "content": query})
            messages.append({"role": "assistant", "content": assistant_response})

            # 构造结构化输出
            structured_output = {
                "user_input": query,
                "model_response": assistant_response,
                "timestamp": datetime.utcnow().isoformat(),
                "messages": messages,
            }
            # 在流式返回的最后，返回结构化输出
            # yield f"{structured_output}"
            a = {"role": "user", "content": query}
            b = {"role": "assistant", "content": assistant_response}
            dataBM.insert_context(user_id, chatbot_id, f"{a}", f"{b}")

            # insert_context(user_id, chatbot_id, f"{b}")

        return generate_response()


class EmbeddingFunc:
    def __init__(self):
        load_dotenv()
        self.api_key = os.getenv('QIANWEN_API_KEY')
        self.client = OpenAI(base_url=f"{LLM_CLOUD_URL}", api_key=self.api_key)

        # openai_ef = OpenAIEmbeddingFunction(
        #     api_key="YOUR_API_KEY",
        #     model_name="text-embedding-nomic-embed-text-v1.5",
        #     api_base=f"{LLM_BASE_URL}/v1"
        # )
        """定义一个OpenAIEmbeddingFunction实例"""
        self.openai_ef = OpenAIEmbeddingFunction(
            api_key=self.api_key,
            model_name="text-embedding-v3",
            api_base=f"{LLM_CLOUD_URL}",
            dimensions=1024,
        )

    """文件处理函数"""

    def process_file(self, filepath: str) -> str:
        ext = os.path.splitext(filepath)[1].lower()

        if ext == '.pdf':
            import fitz  # PyMuPDF 的别名

            text = ""
            with open(filepath, 'rb') as f:
                # 打开 PDF 文件
                pdf_document = fitz.open(f)
                for page_num in range(len(pdf_document)):
                    # 获取每一页的内容
                    page = pdf_document.load_page(page_num)
                    # 提取文本
                    text += page.get_text()

        elif ext == '.docx':
            from docx import Document
            doc = Document(filepath)
            text = '\n'.join([para.text for para in doc.paragraphs])

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    # 处理表格中的文本
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text += cell.text + "\t"
                            text += "\n"

        elif ext == '.xlsx':
            from openpyxl import load_workbook
            wb = load_workbook(filepath, read_only=True)
            text = ""
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    text += "\t".join([str(cell) for cell in row]) + "\n"

        elif ext == '.md':
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()

        elif ext == '.html':
            with open(filepath, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                text = soup.get_text(separator='\n')



        elif ext == '.csv':
            with open(filepath, 'r', newline='', encoding='utf-8') as f:
                reader = csv.reader(f)
                text = "\n".join([",".join(row) for row in reader])

        elif ext == '.txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()

        else:
            raise ValueError(f"Unsupported file format: {ext}")

        # 统一清理文本格式
        text = text.replace('\n', '').replace('\r', '\n').replace('\t', ' ').replace(' ','')
        return text.strip()

    """文本分块函数"""

    import logging

    def chunk_text(self, text: str, max_tokens: int = 100) -> list[str]:
        """Split text into chunks where each chunk has at most `max_tokens` tokens (approximate).

        Args:
            text: Input text to split.
            max_tokens: Maximum tokens per chunk (including spaces).

        Returns:
            List of text chunks.
        """
        if not text.strip():
            return []

        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0  # Track length to avoid repeated joins

        for word in words:
            word_length = len(word) + (1 if current_chunk else 0)  # Add space if not first word
            if current_length + word_length > max_tokens:
                chunks.append(' '.join(current_chunk))
                current_chunk = [word]
                current_length = len(word)
            else:
                current_chunk.append(word)
                current_length += word_length

        if current_chunk:
            chunks.append(' '.join(current_chunk))

        logging.info(f"分割为 {len(chunks)} 个 chunks，首段预览: {chunks[0][:50]}...")
        return chunks

    def split_text(
            self,
            text: str,
            chunk_size: int = 1000,
            chunk_overlap: int = 100,
            recursive: bool = False
    ) -> List[str]:
        """
        将文档分割成较小的文本块

        Args:
            text: 要分割的文本
            chunk_size: 每个文本块的最大字符数
            chunk_overlap: 相邻文本块之间的重叠字符数
            recursive: 是否使用递归分割器

        Returns:
            分割后的文本块列表
        """

        if not text:
            return []

        if recursive:
            return self._recursive_split(text, chunk_size, chunk_overlap)
        else:
            return self._simple_split(text, chunk_size, chunk_overlap)

    def _simple_split(self, text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """简单固定长度分割"""
        logging.info(f"----文件长度{len(text)}")
        if chunk_overlap >= chunk_size:
            raise ValueError("chunk_overlap 必须小于 chunk_size")

        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = min(start + chunk_size, text_length)
            chunks.append(text[start:end])

            # 计算下一个起始位置（考虑重叠）
            start = end - chunk_overlap if end < text_length else end
        logging.info(f"分割为 {len(chunks)} 个 chunks，首段预览: {chunks[0][:50]}...")
        return chunks

    def _recursive_split(self, text: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """递归分割（优先按段落、句子、短语分割）"""
        # 定义分割符优先级（从大到小）
        separators = [
            "\n\n",  # 双换行（段落分隔）
            "\n",  # 单换行
            "(?<=。)",  # 中文句号后
            "(?<=！)",  # 中文感叹号后
            "(?<=？)",  # 中文问号后
            "(?<=;)",  # 分号后
            "(?<=,)",  # 逗号后
            " ",  # 空格
            ""  # 最后按字符分割
        ]

        # 递归分割函数
        def _split_recursively(text: str, separators: List[str]) -> List[str]:
            if len(text) <= chunk_size:
                return [text]

            current_sep = separators[0]
            remaining_seps = separators[1:]

            if current_sep:  # 非空分隔符
                parts = re.split(f"({current_sep})", text)
                # 重新组合分隔符和内容
                combined = []
                for i in range(0, len(parts) - 1, 2):
                    combined.append(parts[i] + parts[i + 1])
                if len(parts) % 2 == 1:
                    combined.append(parts[-1])
            else:  # 空分隔符表示按字符分割
                combined = list(text)

            # 尝试合并小片段
            chunks = []
            current_chunk = ""
            for part in combined:
                if len(current_chunk) + len(part) <= chunk_size:
                    current_chunk += part
                else:
                    if current_chunk:
                        chunks.append(current_chunk)
                    current_chunk = part

            if current_chunk:
                chunks.append(current_chunk)

            # 如果分割后仍有超长块，则使用下一个分隔符递归处理
            if any(len(chunk) > chunk_size for chunk in chunks) and remaining_seps:
                new_chunks = []
                for chunk in chunks:
                    if len(chunk) > chunk_size:
                        new_chunks.extend(_split_recursively(chunk, remaining_seps))
                    else:
                        new_chunks.append(chunk)
                chunks = new_chunks

            return chunks

        chunks = _split_recursively(text, separators)

        # 添加重叠部分
        if chunk_overlap > 0:
            overlapped_chunks = []
            for i in range(len(chunks)):
                if i == 0:
                    overlapped_chunks.append(chunks[i])
                else:
                    prev_end = max(0, len(chunks[i - 1]) - chunk_overlap)
                    overlapped = chunks[i - 1][prev_end:] + chunks[i]
                    overlapped_chunks.append(overlapped)
            chunks = overlapped_chunks

        return chunks




    # OpenAI 分析函数------------------
    def analyze_with_gpt(self, content: str, prompt: str) -> str:
        load_dotenv()
        api_key = os.getenv('QIANWEN_API_KEY')
        client = OpenAI(base_url=f"{LLM_CLOUD_URL}/v1", api_key=api_key)
        response = client.chat.completions.create(
            # model="qwen2.5-7b-instruct-mlx",
            model=LLM_MODEL,
            messages=[
                {"role": "system", "content": "你是一个专业文档分析助手"},
                {"role": "user", "content": f"{prompt}\n\n文档内容：{content}"}
            ],
            temperature=0.3
        )
        return response.choices[0].message.content

    """返回向量化数据"""

    def embedding(self, text):
        client = OpenAI(
            api_key=os.getenv("QIANWEN_API_KEY"),
            base_url=LLM_CLOUD_URL  # 百炼服务的base_url
        )

        completion = client.embeddings.create(
            model="text-embedding-v3",
            input=text,
            dimensions=2000,
            encoding_format="float"
        )

        return completion.model_dump_json()

    # def embedding(text):
    #     url = f"{LLM_BASE_URL}/v1/embeddings"
    #     payload = {
    #         "input": text,
    #         "model": 'text-embedding-nomic-embed-text-v1.5'
    #     }
    #     response = requests.post(url, json=payload)
    #     if response.status_code == 200:
    #         return response.json()["data"][0]["embedding"]
    #     else:
    #         raise Exception(f"Failed to get embedding: {response.status_code}, {response.text}")

    """向向量数据库中添加数据"""

    def vector_db_add(self, texts, collection_name):
        # 初始化向量数据库
        persist_directory = './vector/chroma2'  # 持久化数据  存放处
        # TODO 这里会出现问题，如果没有创建目录，会显示没有权限，这是一个只读目录
        # 如果目录不存在，则创建它（包括所有必要的父目录）
        if not os.path.exists(persist_directory):
            os.makedirs(persist_directory)


        client = chromadb.PersistentClient(persist_directory)
        logging.info(f'开始上传到向量数据库--\n {collection_name}')
        # 添加了embedding函数 TODO 出错点，需要注意，有些参数需要是字符串形式
        collection = client.get_or_create_collection(name=f"{collection_name}", embedding_function=self.openai_ef)
        # 将文件添加到向量数据库中
        collection.add(
            documents=texts,
            metadatas=[{
                "source": "sample",
                "id": str(uuid.uuid4()),  # 生成唯一的 ID
                "timestamp": datetime.now().isoformat()  # 当前时间戳
            } for _ in texts],
            ids=[f"{uuid.uuid4()}{i}" for i in range(len(texts))]
        )
        logging.info('上传到向量数据库成功')

    """在向量数据库中搜索相似值"""

    def search_in_vector_db(self, query,user_id, chatbot_id ):
        # 初始化向量数据库
        persist_directory = './vector/chroma2'  # 持久化数据存放处

        client = chromadb.PersistentClient(persist_directory)
        collection = client.get_collection(name=f"{user_id}_{chatbot_id}", embedding_function=self.openai_ef)
        # collection = client.get_collection(name=f"1_12", embedding_function=self.openai_ef)

        # 将查询文本转换为嵌入向量
        # query_embedding = embedding(query)

        # 在向量数据库中搜索
        results = collection.query(
            # query_embeddings=[query_embedding],
            query_texts=[query],
            n_results=1,
        )
        logging.info(f"-----{query},\n\n---{results}")
        if results['distances'][0][0] > 0.65:
            logging.info('未找到相关信息')
            return ''
        return results['documents'][0][0]
